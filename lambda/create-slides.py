import json
import os
import logging
import boto3
from datetime import datetime
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile

# ログ設定
logger = logging.getLogger()
logger.setLevel(logging.INFO)

# AWSクライアント（リージョン明示）
s3_client = boto3.client('s3', region_name='us-west-2')

def lambda_handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    経営アドバイスのPowerPointプレゼンテーションを生成するLambda関数
    Bedrockエージェント用のアクショングループ関数
    
    Args:
        event: Bedrockエージェントからのイベント（actionGroup, function, parameters含む）
        context: Lambda実行コンテキスト
    
    Returns:
        Bedrockエージェント形式のレスポンス
    """
    try:
        logger.info(f"Received event: {json.dumps(event)}")
        
        # Bedrockエージェントの必須フィールドを取得
        action_group = event['actionGroup']
        function = event['function']
        message_version = event.get('messageVersion', 1)
        parameters = event.get('parameters', [])
        
        # S3バケット名の取得
        bucket_name = os.environ.get('S3_BUCKET_NAME')
        if not bucket_name:
            raise ValueError("S3_BUCKET_NAME environment variable is not set")
        
        # パラメータから必要なデータを取得
        param_dict = {}
        for param in parameters:
            param_dict[param['name']] = param['value']
        
        analysis_data = param_dict.get('analysis_data', '')
        market_data = param_dict.get('market_data', '')
        sales_data = param_dict.get('sales_data', '')
        title = param_dict.get('title', '鳥豪族 経営アドバイスレポート')
        
        # PowerPointプレゼンテーションの作成
        ppt = create_presentation(analysis_data, market_data, sales_data, title)
        
        # 一時ファイルに保存
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"toriguozoku_report_{timestamp}.pptx"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            ppt.save(tmp_file.name)
            tmp_file_path = tmp_file.name
        
        # S3にアップロード
        s3_key = f"presentations/{filename}"
        s3_client.upload_file(tmp_file_path, bucket_name, s3_key)
        
        # 一時ファイルを削除
        os.unlink(tmp_file_path)
        
        # 署名付きURLの生成（24時間有効）
        try:
            presigned_url = s3_client.generate_presigned_url(
                'get_object',
                Params={'Bucket': bucket_name, 'Key': s3_key},
                ExpiresIn=86400  # 24時間
            )
            logger.info(f"署名付きURL生成成功: {presigned_url}")
        except Exception as url_error:
            logger.error(f"署名付きURL生成エラー: {str(url_error)}")
            presigned_url = f"https://{bucket_name}.s3.us-west-2.amazonaws.com/{s3_key}"
        
        logger.info(f"PowerPoint生成完了: {s3_key}")
        
        # Bedrockエージェント形式のレスポンス
        response_body = {
            'TEXT': {
                'body': json.dumps({
                    'message': 'PowerPointプレゼンテーションが正常に生成されました',
                    'download_url': presigned_url,
                    'filename': filename,
                    's3_key': s3_key
                }, ensure_ascii=False)
            }
        }
        
        action_response = {
            'actionGroup': action_group,
            'function': function,
            'functionResponse': {
                'responseBody': response_body
            }
        }
        
        response = {
            'response': action_response,
            'messageVersion': message_version
        }
        
        logger.info('Response: %s', response)
        return response
        
    except KeyError as e:
        logger.error('必須フィールドが不足しています: %s', str(e))
        return {
            'statusCode': 400,
            'body': f'エラー: 必須フィールドが不足しています - {str(e)}'
        }
    except Exception as e:
        logger.error(f"エラーが発生しました: {str(e)}")
        return {
            'statusCode': 500,
            'body': f'PowerPoint生成でエラーが発生しました: {str(e)}'
        }

def create_presentation(analysis_data: str, market_data: str, sales_data: str, title: str) -> Presentation:
    """
    PowerPointプレゼンテーションを作成
    
    Args:
        analysis_data: 分析データ（文字列）
        market_data: 市場データ（文字列）
        sales_data: 売上データ（文字列）
        title: プレゼンテーションタイトル
    
    Returns:
        作成されたPresentationオブジェクト
    """
    prs = Presentation()
    
    # スライド1: タイトルスライド
    create_title_slide(prs, title)
    
    # スライド2: エグゼクティブサマリー
    create_executive_summary_slide(prs, analysis_data)
    
    # スライド3: 売上分析
    create_sales_analysis_slide(prs, sales_data)
    
    # スライド4: 市場動向
    create_market_trends_slide(prs, market_data)
    
    # スライド5: 改善提案
    create_recommendations_slide(prs, analysis_data)
    
    # スライド6: アクションプラン
    create_action_plan_slide(prs, analysis_data)
    
    return prs

def create_title_slide(prs: Presentation, title: str):
    """タイトルスライドを作成"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Generated on {datetime.now().strftime('%Y年%m月%d日')}"
    
    # タイトルの書式設定
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

def create_executive_summary_slide(prs: Presentation, analysis_data: Dict):
    """エグゼクティブサマリースライドを作成"""
    slide_layout = prs.slide_layouts[1]  # タイトル + コンテンツレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "エグゼクティブサマリー"
    
    # サマリー内容
    summary_points = [
        "• 売上トレンド分析結果",
        "• 主要商品パフォーマンス",
        "• 市場機会の特定",
        "• 競合他社との比較分析",
        "• 収益改善の重点施策"
    ]
    
    content_shape.text = "\n".join(summary_points)
    
    # 書式設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)

def create_sales_analysis_slide(prs: Presentation, sales_data: str):
    """売上分析スライドを作成"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "売上分析"
    
    # 売上データの表示
    sales_points = [
        "📊 売上パフォーマンス分析",
        f"• 売上データ:",
        sales_data,
        "",
        "🎯 改善ポイント",
        "• 低パフォーマンス商品の見直し",
        "• 地域別戦略の最適化",
        "• 季節要因の活用"
    ]
    
    content_shape.text = "\n".join(sales_points)
    
    # 書式設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(16)

def create_market_trends_slide(prs: Presentation, market_data: str):
    """市場動向スライドを作成"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "市場動向分析"
    
    # 市場データの表示
    market_points = [
        "🌟 市場データ",
        market_data,
        "",
        "💡 市場機会",
        "• 新規出店エリアの特定",
        "• 限定メニューの展開可能性",
        "• 差別化ポイントの明確化が重要"
    ]
    
    content_shape.text = "\n".join(market_points)
    
    # 書式設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(16)

def create_recommendations_slide(prs: Presentation, analysis_data: str):
    """改善提案スライドを作成"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "改善提案"
    
    # 提案内容
    recommendations = [
        "📋 分析結果・改善提案",
        analysis_data,
        "",
        "📈 実施スケジュール",
        "• 短期施策（1-3ヶ月）",
        "• 中期施策（3-6ヶ月）", 
        "• 長期施策（6ヶ月以上）"
    ]
    
    content_shape.text = "\n".join(recommendations)
    
    # 書式設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(16)

def create_action_plan_slide(prs: Presentation, analysis_data: str):
    """アクションプランスライドを作成"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "アクションプラン"
    
    # アクションプラン
    action_items = [
        "📅 今後30日間の重点取り組み",
        "",
        "週次レビュー項目:",
        "• 売上実績の確認と分析",
        "• 顧客満足度調査の実施",
        "• スタッフフィードバックの収集",
        "",
        "月次改善サイクル:",
        "• KPI達成状況の評価",
        "• 施策効果の測定",
        "• 次月計画の策定",
        "",
        "📊 成功指標 (KPI)",
        "• 月間売上成長率: +10%目標",
        "• 顧客リピート率: 向上",
        "• 新規顧客獲得数: 増加"
    ]
    
    content_shape.text = "\n".join(action_items)
    
    # 書式設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(16)